# -*- coding: utf-8 -*-
from __future__ import annotations

import base64
import io
import os
import re
from pathlib import Path


def _format_usage_table(usage_log: list[dict]) -> str:
    if not usage_log:
        return "No usage recorded."
    totals: dict[str, dict] = {}
    for entry in usage_log:
        t = totals.setdefault(entry["label"], {"calls": 0, "prompt": 0, "completion": 0})
        t["calls"] += 1
        t["prompt"] += entry["prompt_tokens"]
        t["completion"] += entry["completion_tokens"]
    lines = [f"{'Label':<24}{'Calls':>7}{'Prompt':>10}{'Completion':>12}{'Total':>10}"]
    grand_total = 0
    for label, t in sorted(totals.items(), key=lambda kv: -(kv[1]["prompt"] + kv[1]["completion"])):
        total = t["prompt"] + t["completion"]
        grand_total += total
        lines.append(f"{label:<24}{t['calls']:>7}{t['prompt']:>10}{t['completion']:>12}{total:>10}")
    lines.append("-" * 63)
    lines.append(f"{'TOTAL':<24}{'':>7}{'':>10}{'':>12}{grand_total:>10}")
    return "\n".join(lines)


class MetadataExtractor:

    # The first-pass tile target is deliberately a little larger than the old 1900 px
    # setting. OpenAI high-detail vision downsamples large crops internally, so 1900 px
    # caused many more requests without a proportional gain in effective model resolution.
    # Hard/ambiguous areas are still re-cropped at higher magnification below.
    BASE_TILE_TARGET = 2250

    def __init__(self):
        # In-memory cache prevents a late downstream failure from triggering the entire
        # expensive vision pipeline a second time in FileAgent's fallback path.
        self._vision_cache: dict[tuple[str, int, int, str], str] = {}
        self.usage_log: list[dict] = []

    def _track_usage(self, label: str, usage) -> None:
        if usage is None:
            return
        self.usage_log.append({
            "label": label,
            "prompt_tokens": getattr(usage, "prompt_tokens", 0) or 0,
            "completion_tokens": getattr(usage, "completion_tokens", 0) or 0,
            "total_tokens": getattr(usage, "total_tokens", 0) or 0,
        })

    def usage_summary(self) -> str:
        return _format_usage_table(self.usage_log)

    VALID_ISA_PREFIXES = {
        "PI", "PDI", "PDIT", "PT", "PIT", "PG", "PCV", "PSV", "PSE", "PAH", "PAL",
        "PAHH", "PALL", "PIC", "PR", "PY", "PDT",
        "TI", "TT", "TE", "TC", "TIC", "TR", "TCV", "TAH", "TAL",
        "FI", "FT", "FE", "FIT", "FC", "FCV", "FIC", "FR", "FAH", "FAL", "FL",
        "LI", "LT", "LE", "LC", "LCV", "LIC", "LAH", "LAL", "LAHH", "LALL",
        "AI", "AT", "AE", "AC", "ACV", "AIC",
        "SI", "ST", "SE", "SC", "SCV", "SIC", "SV",
        "VI", "VT", "VE", "VC",
        "HV", "XV", "ZV", "YV",
        "HS", "XS", "ZS", "YS",
        "HI", "XI", "ZI", "YI",
        "PH", "PL", "TH", "TL", "FH", "FL", "LH", "LL",
        "PAH", "TAH", "FAH", "LAH",
    }

    TILE_INSTRUMENT_PROMPT = (
        "You are an ISA 5.1 instrumentation specialist reading a CROPPED SECTION of a P&ID or engineering diagram.\n\n"
        "YOUR ONLY JOB: List every instrument tag bubble visible in this image section.\n\n"
        "STRICT RULES:\n"
        "- Instrument tags ONLY appear inside circles, ovals, squares, or hexagons drawn on the diagram\n"
        "- Valid ISA prefixes: PI, PDI, PDIT, PDT, PT, PIT, PG, PCV, PSV, PSE, PAH, PAL, PAHH, PALL, PH, "
        "TI, TT, TE, TC, TCV, TAH, TAL, "
        "FI, FT, FE, FIT, FC, FCV, FAH, FAL, FL, "
        "LI, LT, LE, LC, LCV, LAH, LAL, LAHH, LALL, "
        "SV, SI, ST, SE, HV, XV, ZV, HS, XS, ZS, "
        "PIT, FIC, LIC, TIC, PIC\n"
        "- DO NOT read numbers from: title blocks, revision tables, coordinates, borders, dates, "
        "P.O. numbers, contract numbers, drawing numbers, personnel initials/signatures\n"
        "- CHARACTER AMBIGUITY -- on engineering drawings these pairs are frequently confused by OCR. "
        "Always resolve in favour of the ISA-valid interpretation:\n"
        "  * 0 vs O vs C: a circle-like character in a tag prefix is almost always 0 (zero), not O or C\n"
        "  * 1 vs I vs l: a vertical stroke in a tag prefix is almost always I (the ISA letter), not 1 or l\n"
        "  * Example: if you see 'C4 PDI' or '4C PDI' written on equipment, the likely reading is '44 PDI' (a two-digit equipment number, then the ISA prefix)\n"
        "  * Example: if you see PD1 in a bubble, the likely correct reading is PDI (Pressure Differential Indicating)\n"
        "- COMPOUND BUBBLES: some instrument bubbles contain TWO stacked tags in a single circle "
        "(e.g. PDI on top and PDIT below in the same bubble, or FIT above FI). "
        "These are TWO separate instruments sharing a bubble -- list BOTH on separate lines\n"
        "- ASSOCIATED TRANSMITTERS: wherever you see a PDI, also look for a co-located PDIT. "
        "Wherever you see an FI, look for a co-located FIT. List each separately if present\n"
        "- If you cannot clearly read a tag, write UNREADABLE\n"
        "- DO NOT guess or fabricate. If unsure of a digit, write the prefix and UNREADABLE (e.g. PDI-????)\n"
        "- DO NOT FABRICATE A TAG BY COMBINING UNRELATED ELEMENTS: a valid tag's prefix and number "
        "are drawn together, as ONE unit, inside or directly on a single instrument bubble. If you "
        "see a bare 2-4 digit number that is NOT inside an instrument bubble (e.g. it's a BOM "
        "reference number sitting next to a valve or fitting symbol), do NOT attach a nearby, "
        "separate ISA prefix letter to it just because both appear somewhere in this crop -- they "
        "belong to two different things. Only report a tag when the prefix letters and the number "
        "are drawn as a single unit inside one bubble.\n"
        "- Note any HI/LO/HH/LL setpoint values shown adjacent to bubbles\n\n"
        "OUTPUT FORMAT -- return ONLY a plain list, one tag per line, nothing else:\n"
        "PDI-XXXX (HI, LO)\n"
        "PDIT-XXXX\n"
        "FIT-YYYY (HI, LO)\n"
        "FCV-YYYY\n"
        "UNREADABLE\n"
    )

    CONTEXT_PROMPT = (
        "You are a senior process engineer reading a full engineering diagram.\n\n"
        "Extract ONLY the following -- be factual, do not guess:\n\n"
        "CHARACTER AMBIGUITY NOTE: On engineering drawings, 0 (zero) and O, and 1 (one) and I (capital i) "
        "are frequently confused. When reading equipment labels: prefer numeric reading for standalone numbers "
        "(e.g. '70 PDI COMPRESSOR' not '7C PDI' or '70 PD1'); prefer letter reading for ISA prefixes. "
        "Read equipment labels exactly as they appear character by character and favour the reading that "
        "makes engineering sense (e.g. a compressor labelled '70 PDI' is a 70-PDI type unit).\n\n"
        "1. DOCUMENT TYPE: (P&ID, PFD, System Diagram, GA, Isometric, etc.)\n\n"
        "2. TITLE BLOCK: Extract exactly as written -- preserve all distinctions:\n"
        "   - Drawing title (e.g. 'GAS SEAL SYSTEM DIAGRAM')\n"
        "   - Drawing number / document number (this is the vendor or engineering firm's internal drawing ID, "
        "NOT the project number -- label it clearly as DRAWING NUMBER)\n"
        "   - Revision number and date\n"
        "   - Vendor / contractor name (e.g. DRESSER-RAND, KBR, Worley)\n"
        "   - Client name (e.g. TRANSCONTINENTAL GAS PIPE LINE)\n"
        "   - Project name (e.g. MID ATLANTIC CONNECTOR EXPANSION PROJECT) -- "
        "this is NOT the same as the drawing number; label it clearly as PROJECT NAME\n"
        "   - Work order / purchase order number if present -- label as W.O. or P.O., NOT as project number\n"
        "   - Contract number if present\n"
        "   - Sheet number (e.g. SH. 1 OF 2)\n"
        "   - IMPORTANT: Do NOT use the drawing number or W.O./P.O. number as the project number. "
        "If no explicit project number is visible, write PROJECT NUMBER: NOT STATED\n\n"
        "3. MAJOR EQUIPMENT: List each piece of process equipment with its exact label as written:\n"
        "   - Compressors, pumps, motors, drivers -- read labels carefully, distinguishing 0 from O/C and 1 from I\n"
        "   - Vessels, drums, tanks\n"
        "   - Heat exchangers, coolers\n"
        "   - Skid boundaries and panel boundaries (dashed box labels)\n"
        "   A piece of equipment must have its OWN drawn symbol (a vessel outline, a compressor "
        "shape, a shell, a named box) and its own name/label physically on or immediately next "
        "to that symbol. A pipe spec string (pattern SIZE-CLASS-SUFFIX) is NEVER itself a piece "
        "of equipment -- if you only see a spec-like string with no separate drawn equipment "
        "symbol and no equipment name attached to it, do not invent an equipment entry from it "
        "(e.g. do not attach a type like 'HEAT EXCHANGER' or 'VESSEL' to a string that is "
        "actually a pipe spec).\n\n"
        "4. PIPE SPECIFICATIONS: Do NOT attempt to precisely transcribe pipe spec label "
        "characters here -- a dedicated higher-magnification pass elsewhere in this pipeline "
        "handles that specifically, with safeguards for common character confusion (B/8, D/0, "
        "etc.) that this pass doesn't have. Just note roughly how many distinct pipe spec "
        "labels appear to be present, for orientation only.\n\n"
        "5. PROCESS STREAMS AND CONNECTIVITY: List EVERY named stream visible, not just the "
        "most prominent one -- this specifically includes ALTERNATE, BACKUP, SECONDARY, or "
        "SPARE supply/vent/drain lines, which are easy to overlook next to a main line of the "
        "same kind but are frequently drawn and labelled separately. Supply lines, vent lines, "
        "drain lines. For each named stream, also note WHERE it enters or exits the diagram (a "
        "labelled connection point, arrow, or edge-of-drawing tag) and, if visible, what "
        "equipment it connects to first -- this connectivity information is needed downstream "
        "to trace fluid paths, so capture it whenever it's legible even if brief. Describe the "
        "connection by NAMED EQUIPMENT or CONNECTION POINT only (e.g. 'SUPPLY GAS enters at "
        "left edge, goes to gas seal panel') -- do NOT attempt to transcribe the pipe spec "
        "characters here; that is handled by the dedicated, safeguarded pass elsewhere in this "
        "pipeline, and a spec string written here with no cross-checking is a common source of "
        "character-level misreads (e.g. confusing a B for an 8, or a D for a 0) entering the "
        "final output unchecked.\n\n"
        "6. NOTES/SAFETY: Any general notes, safety annotations, or legend content\n\n"
        "7. UNIQUE ELEMENTS: Only note elements that have specific engineering significance: "
        "unusual equipment configurations, non-standard connections, safety-critical markings, "
        "vendor-specific assemblies, or process features that distinguish this diagram from a typical one of its type. "
        "Do NOT note generic observations like 'diagram has a title block' or 'boundary boxes are present'.\n\n"
        "DO NOT list instrument tags here -- those are handled separately.\n"
        "If you cannot read something clearly, omit it rather than guessing.\n"
    )

    RECONCILIATION_PROMPT = (
        "You are an ISA 5.1 instrumentation specialist.\n\n"
        "Below is a RAW LIST of instrument tags extracted from multiple tile scans of a P&ID. "
        "Some may be duplicates, some may be misread, and some may be fabricated by the AI.\n\n"
        "Your job:\n"
        "1. Deduplicate -- merge identical tags\n"
        "2. Correct obvious character misreads using these rules:\n"
        "   - In tag PREFIXES: the letter I (capital i) is almost always an ISA function letter, not the digit 1\n"
        "     e.g. PD1-XXXX should be corrected to PDI-XXXX\n"
        "   - In tag PREFIXES: the digit 0 (zero) should not appear -- if you see it, "
        "check whether it is a misread O or whether the whole prefix is invalid\n"
        "   - In tag NUMBERS (after the hyphen): digits only are expected; letters I and O "
        "should be corrected to 1 and 0 respectively if they appear in the number portion\n"
        "3. Where both PDI-XXXX and PDIT-XXXX appear with the same loop number, keep BOTH -- "
        "they are two separate instruments in the same loop\n"
        "4. Remove any tags whose prefix is not a valid ISA function letter combination\n"
        "5. Remove any tags that look like drawing numbers, dates, P.O. numbers, or contract numbers "
        "(7+ digit numbers, numbers containing slashes, numbers starting with year patterns like 19xx or 20xx)\n"
        "6. Flag any tag you remain uncertain about with a ? suffix\n\n"
        "RAW TAG LIST:\n"
        "{raw_tags}\n\n"
        "Return ONLY a clean deduplicated list, one tag per line. No commentary.\n"
    )

    PIPE_SPEC_PROMPT = (
        "You are reading a NARROW HORIZONTAL STRIP of a P&ID engineering drawing.\n\n"
        "YOUR ONLY JOB: Find and list every pipe specification label visible in this strip.\n"
        "Pipe specs appear as text inside rectangular boxes on or adjacent to pipe lines.\n"
        "They follow patterns like: SIZE-CLASS-SUFFIX (e.g. 0.5-XB2-1, 1.5-YC4-2A, 2.0-ZD6-3C)\n\n"
        "CHARACTER AMBIGUITY -- this is critical for pipe specs:\n"
        "- The letter B (uppercase B) and the digit 8 are commonly confused. "
        "Look at the character carefully: B has two bumps on the right side, 8 has two symmetric loops.\n"
        "- D and 0: D has a flat left vertical stroke, 0 is a closed oval.\n"
        "- Read EVERY character literally as you see it. Do NOT substitute or correct.\n"
        "- If a spec reads 0.5-XB2-7, write 0.5-XB2-7. Do not write 0.5-X82-7 (B/8 are commonly confused).\n"
        "- If a spec reads 0.5-YD4-7, write 0.5-YD4-7. Do not write 0.5-Y04-7 (D/0 are commonly confused).\n\n"
        "UNCERTAINTY -- read this before answering:\n"
        "- If a single character is genuinely unclear even after close inspection, mark that "
        "exact character position with ? (e.g. 1.0-3?6-416C). Do NOT invent a complete, "
        "clean-looking string for a label you cannot fully read -- a partial read with ? is "
        "far more useful than a confident but wrong one.\n"
        "- If two spec labels sit close together, read each one completely and independently "
        "before moving to the next. Do not let digits from one label bleed into your reading "
        "of an adjacent label -- this is a common error when labels are stacked or clustered.\n"
        "- If you are not looking at a real spec label at all (e.g. it's a dimension, a BOM "
        "reference number, or unreadable noise), do not report it as a pipe spec.\n\n"
        "There is also a NOTE BOX typically in the lower left area that says:\n"
        "  [pipe spec] TYP. FOR INSTRUMENTATION AND REFERENCE SIGNAL LINES\n"
        "Read that spec carefully -- it defines the typical spec for instrument lines. "
        "If you find this note box, prefix that one line with 'TYP:' (e.g. 'TYP: 0.5-XB2-1'). "
        "All other specs get their own plain line, no prefix.\n\n"
        "OUTPUT: one pipe spec per line, nothing else.\n"
        "Example:\n"
        "TYP: 0.5-XB2-1\n"
        "1.0-YC4-2\n"
        "2.0-ZD6-3C\n"
    )

    PIPE_SPEC_DETAIL_ZOOM_PROMPT = (
        "You are re-examining a SMALL, HIGH-MAGNIFICATION crop containing one or more pipe "
        "specification labels that were flagged as uncertain or densely clustered on a first "
        "pass.\n\n"
        "RULES:\n"
        "- Read each spec label completely and independently before moving to the next. Do "
        "not let digits from one nearby label influence your reading of another -- this is "
        "the most common error in dense clusters.\n"
        "- Character ambiguity: B vs 8 (B has two flat-sided bumps, 8 has two symmetric "
        "loops), D vs 0 (D has a flat left edge, 0 is a closed oval). Read exactly as drawn.\n"
        "- If a character is genuinely illegible even at this magnification, mark that "
        "position with ? -- do not guess a complete, clean-looking string.\n"
        "- If what looked like one label on the first pass is actually two adjacent distinct "
        "labels, report them as two separate lines.\n\n"
        "OUTPUT: one pipe spec per line, nothing else. Use ? for any illegible character.\n"
    )

    PIPE_SPEC_RECONCILE_PROMPT = (
        "You are a piping engineer reconciling pipe specification readings extracted from "
        "multiple overlapping crops of the same P&ID.\n\n"
        "Below is the RAW LIST of pipe spec strings as read, in the order encountered. The "
        "same physical spec label may have been read more than once (from overlapping crops "
        "or a follow-up zoom pass) with slightly different results due to vision noise. Some "
        "entries may be hallucinated -- plausible-looking strings assembled from digits of two "
        "adjacent real labels, or invented outright when a crop was blurry or ambiguous.\n\n"
        "A confirmed reference spec (if the drawing's 'TYP. FOR INSTRUMENTATION' note box was "
        "found) is provided separately -- treat it as correct and cross-check other entries "
        "against it.\n\n"
        "YOUR JOB:\n"
        "1. Group near-identical strings that differ by 1-2 characters -- these are almost "
        "always the same physical label read with minor noise. Pick the single most plausible "
        "canonical reading per group: prefer whichever variant appears more than once, or "
        "whichever matches the confirmed reference spec if related.\n"
        "2. A spec string that appears only ONCE across all readings AND is not corroborated "
        "by the reference spec AND looks structurally suspicious (e.g. its digits look like a "
        "blend of two other specs in the list) is LOW CONFIDENCE. Do not silently keep or "
        "silently discard it -- list it separately so a human can verify against the drawing.\n"
        "3. Do not invent a spec that doesn't appear anywhere in the raw list.\n\n"
        "REFERENCE NOTE SPEC (if any): {reference_spec}\n\n"
        "RAW SPEC READINGS:\n{raw_specs}\n\n"
        "OUTPUT FORMAT:\n"
        "CONFIRMED SPECS:\n- <canonical spec string>\n\n"
        "LOW CONFIDENCE (verify against drawing):\n- <spec string>: <one clause on why it's suspect>\n"
    )

    VALVE_SURVEY_PROMPT = (
        "You are an engineering drawing specialist examining a CROP of a P&ID.\n\n"
        "YOUR ONLY JOB: Catalogue every valve, fitting, and piping specialty item visible.\n\n"
        "VALVE SYMBOL GUIDE for P&IDs:\n"
        "- Check valve: bow-tie or arrowhead symbol on a line (allows flow one direction only)\n"
        "- Ball valve: filled solid square or circle on a line (quarter-turn isolation)\n"
        "- Gate valve: filled solid triangle pointing at line (on/off isolation)\n"
        "- Globe valve: circle with a line through it, or bowtie with circle (control/throttling)\n"
        "- Needle valve: small triangle or X symbol (fine flow control, instrumentation)\n"
        "- Control valve: globe/rotary body with actuator on top (diaphragm dome = pneumatic)\n"
        "- Solenoid valve: square symbol with S inside or coil symbol\n"
        "- Rupture disc / PSE: parallel lines or rectangular block across a line (bursting disc)\n"
        "- Orifice plate / restriction: double vertical lines across a pipe (flow restriction)\n"
        "- Filter element: diamond shape (FL- tagged)\n"
        "- Strainer: Y-shape or basket symbol\n"
        "- Spectacle blind: figure-8 symbol\n"
        "- Drain: small line with arrow pointing down\n"
        "- Vent: small line with arrow pointing up\n\n"
        "RULES:\n"
        "- List EVERY fitting visible regardless of size\n"
        "- For each item: identify its type, the pipe spec label on the line it is on, "
        "any tag (e.g. FCV-XXXX, FL-XXXXA), and any adjacent reference number\n"
        "- Reference numbers are BARE 2-4 DIGIT NUMBERS (no letters) near fittings -- "
        "note them as BOM reference numbers; state that a project BOM is needed to confirm meaning\n"
        "- DO NOT confuse a BOM reference number with an OFF-PAGE CONTINUATION label. Small "
        "circled labels made of ONE LETTER plus ONE DIGIT (e.g. K3, L3, U3, Y3), typically found "
        "where a line exits the drawing boundary or connects to another sheet, indicate a "
        "continuation to/from another drawing -- they are NOT a fitting's BOM reference number "
        "and do NOT have a valve/fitting type. Do not attach a type to one of these (in "
        "particular, do not read the letter Y in a continuation label like 'Y3' as if it were "
        "the Y-shaped strainer symbol described above -- a continuation label is text in a "
        "circle, not a drawn strainer shape). If you see one, ignore it for this survey.\n"
        "- Do NOT read instrument bubble tags (circles with text inside) -- focus only on fittings\n"
        "- If you cannot determine valve type, write UNKNOWN\n"
        "- A dashed rectangle around a group of fittings indicates a typical/repeated assembly -- "
        "treat any crop containing part of a dashed boundary box as a package/panel boundary "
        "that likely packs several small items close together. Be extra thorough here: list "
        "every distinct symbol individually, even ones partially cut off at the tile edge, "
        "rather than summarising the group.\n"
        "- SPATIAL ATTRIBUTION: When several reference numbers are stacked or clustered close "
        "together near a group of equipment (e.g. a filter with isolation valves on either "
        "side, or several small symbols packed into one area), each number belongs to its OWN "
        "nearest distinct symbol -- work through the cluster one number at a time, matching "
        "each to the single symbol it is closest to or connected to by a leader line. Do not "
        "assign the same reference number to two different equipment types (e.g. a valve AND "
        "the filter beside it) just because they sit near each other -- a stack of several "
        "numbers almost always labels several DIFFERENT nearby items, not one item repeated.\n"
        "- GENERIC/TYPICAL reference numbers: if a number sits inside or near a box labelled "
        "'TYP.' or 'TYPICAL', it identifies a SYMBOL TYPE that legitimately recurs at multiple "
        "separate locations on the drawing, not one unique item. Report each occurrence where "
        "you find it, with its own correctly-read type and line -- do not treat the same "
        "number appearing more than once, in different places, as an error to resolve.\n\n"
        "OUTPUT FORMAT -- one item per line:\n"
        "- <valve/fitting type> | Line: <pipe spec> | Tag: <if present> | Ref#: <BOM number if present>\n"
        "Example:\n"
        "- Check valve | Line: 0.5-XB2-1 | Tag: none | Ref#: 101 (BOM ref -- project BOM needed)\n"
        "- Ball valve | Line: 0.5-XB2-1 | Tag: none | Ref#: 102 (BOM ref -- project BOM needed)\n"
        "- Rupture disc | Line: 2.0-ZD6-3C | Tag: none | Ref#: 103 (BOM ref -- project BOM needed)\n"
        "- Orifice plate | Line: 0.5-XB2-1 | Tag: none | Ref#: 104\n"
    )

    # --- Verification pass: re-examines a tile flagged as ambiguous by the first valve
    # survey pass. Distinct from PIPE_SPEC_PROMPT (previously aliased here in error --
    # that alias meant the "detail zoom" pass, when it existed, was reading pipe specs
    # a second time instead of resolving symbol ambiguity).
    VALVE_DETAIL_ZOOM_PROMPT = (
        "You are re-examining a SMALL, HIGH-MAGNIFICATION crop of a P&ID or diagram. "
        "This crop was flagged as ambiguous or symbol-dense on a first pass -- your job is "
        "to resolve it with certainty, not to re-survey the whole area.\n\n"
        "SYMBOL DISAMBIGUATION -- read this before answering:\n"
        "- GATE valve vs NEEDLE valve: both can render as a triangle. A GATE valve's triangle "
        "sits directly ON the main process line, sized to match the line weight, no offset. "
        "A NEEDLE valve's triangle is SMALL, sits on a thin instrument impulse/takeoff line "
        "branching off the main line, and is usually near an instrument bubble it feeds. "
        "If the triangle is on a thin branch line near an instrument tag, it is a needle valve.\n"
        "- CHECK valve vs BALL valve vs GATE valve: a check valve is a bow-tie / two opposing "
        "triangle arrowhead shape indicating one-way flow -- it has NO separate operating "
        "handle or actuator stem drawn. A ball valve is typically a single filled circle or "
        "square directly on the line, often with a short perpendicular stem for the handle. A "
        "gate valve is a single solid triangle with a stem/handwheel on top. If the symbol is "
        "two triangles point-to-point (or a bow-tie outline) with NO stem drawn, it is a check "
        "valve -- do not classify it as ball or gate just because it's roughly similar in size.\n"
        "- ORIFICE PLATE / RESTRICTION vs CHECK VALVE: an orifice/restriction is two parallel "
        "bars perpendicular to the line, no directional arrowhead. An arrowhead or wedge shape "
        "means check valve, not orifice.\n"
        "- Reference numbers (2-4 digit numbers near a symbol) identify a BOM line item, not a "
        "valve type -- never infer type from the number, only from the drawn symbol shape.\n"
        "- DO NOT confuse a BOM reference number with an OFF-PAGE CONTINUATION label. A small "
        "circled ONE LETTER + ONE DIGIT label (e.g. K3, U3, Y3) marks a continuation to another "
        "drawing sheet, not a fitting reference -- it has no valve/fitting type, ignore it here.\n"
        "- SPATIAL ATTRIBUTION: if this crop shows multiple reference numbers stacked or "
        "clustered near one piece of equipment, match each number to its own single nearest "
        "symbol individually -- do not let two numbers collapse onto one symbol, and do not "
        "let one symbol's number bleed onto a neighboring symbol. Work through the cluster "
        "top-to-bottom or left-to-right, one number at a time.\n"
        "- If a number sits inside or near a box labelled 'TYP.'/'TYPICAL', it labels a symbol "
        "TYPE reused at multiple separate locations, not a single unique item -- report what "
        "you see at THIS location on its own merits, without treating a repeated number as "
        "evidence you've mis-read something.\n\n"
        "For EACH symbol in this crop, give your answer AND a short reason distinguishing it "
        "from the confusable alternative above.\n\n"
        "OUTPUT FORMAT -- one item per line:\n"
        "- <valve/fitting type> | Line: <pipe spec if visible> | Tag: <if present> | "
        "Ref#: <if present> | Why: <short clause>\n"
        "If, after this close look, you are still genuinely uncertain, write "
        "CANDIDATES: <type A> or <type B> instead of guessing a single type.\n"
    )

    VALVE_RECONCILE_PROMPT = (
        "You are a piping engineer reconciling a valve/fitting survey assembled from multiple "
        "overlapping crops of the same P&ID, plus the separately-extracted instrument tag list "
        "from the same drawing. Each block of raw readings below is labelled with the tile it "
        "came from, e.g. '[Tile R2C3]' -- this tells you WHERE on the drawing that reading "
        "came from.\n\n"
        "CRITICAL RULE #1 -- READ BEFORE ANYTHING ELSE:\n"
        "An item identified by an explicit alphanumeric instrument tag (e.g. PSE-XXXX, "
        "PDI-YYYY -- letters followed by a hyphen and numbers) is a DIFFERENT PHYSICAL DEVICE "
        "from an item identified only by a bare 2-4 digit BOM reference number (e.g. 501), "
        "even if they are the same device TYPE (e.g. both rupture discs) and appear in the "
        "same general area of the drawing. NEVER merge a tagged instrument with an untagged "
        "BOM-ref fitting. If both appear in the source material, both must appear separately "
        "in your output, each keeping its own identifier.\n\n"
        "CRITICAL RULE #2 -- USE TILE POSITION, NOT JUST THE NUMBER, TO DECIDE WHAT'S A DUPLICATE:\n"
        "P&ID reference numbers come in two kinds, and confusing them causes real errors:\n"
        "  (a) UNIQUE instance identifiers -- refer to exactly one physical component.\n"
        "  (b) GENERIC/TYPICAL callouts (often near a box marked 'TYP.') -- the SAME small "
        "number legitimately labels EVERY occurrence of a recurring symbol type, at MULTIPLE "
        "separate physical locations on the drawing. This is normal, not an error.\n"
        "You cannot tell which kind a number is from the number alone -- use tile position:\n"
        "  - If the SAME ref#/tag appears in the SAME tile or an ADJACENT tile (row and column "
        "each differ by at most 1, e.g. R2C3 and R2C4), it is almost certainly ONE physical "
        "item seen twice because tile crops overlap. Merge into one clean entry.\n"
        "  - If the SAME ref#/tag appears in tiles that are NOT adjacent (e.g. R2C3 and R5C8), "
        "it is almost certainly TWO DIFFERENT physical items that happen to share a reused "
        "generic/typical label. Keep them as SEPARATE entries, each with its own tile and its "
        "own independently-read type -- do NOT merge them and do NOT flag them as conflicting, "
        "since there is no actual conflict: they are different real items.\n\n"
        "YOUR JOB:\n"
        "1. Apply Rule #2 above to decide, for every repeated ref#/tag, whether it's one item "
        "re-seen (merge) or several distinct items sharing a label (keep separate).\n"
        "2. When merging a genuine re-seen item, prefer the zoom-verification reading over the "
        "general survey reading if both exist for it.\n"
        "3. Do NOT merge two entries just because they share a type and sit near each other in "
        "the text -- only merge per Rule #2.\n"
        "4. Preserve every distinct item as its own entry. Do not drop an item just because "
        "it's the only mention of that ref# in its tile.\n"
        "5. If the SAME tile (or adjacent tiles) reports CONFLICTING types for what is clearly "
        "one physical item (e.g. 'check valve' in R2C3, 'gate valve' in R2C4, same ref#, "
        "immediately adjacent), that IS a genuine reading conflict -- list it as UNRESOLVED.\n\n"
        "INSTRUMENT TAGS ON THIS DRAWING (for cross-reference only -- do not re-list these as "
        "valve/fitting entries):\n{instrument_tags}\n\n"
        "RAW VALVE/FITTING SURVEY (tile-labelled blocks from multiple overlapping crops, may "
        "include zoom-verification readings):\n{raw_survey}\n\n"
        "OUTPUT FORMAT -- one item per line, keep the Tile field so the position that "
        "justified your merge/keep-separate decision stays visible:\n"
        "- <valve/fitting type> | Line: <pipe spec if known> | Tag: <if present> | "
        "Ref#: <if present> | Tile: <the tile this reading is from, or the merged tile if "
        "combined>\n\n"
        "UNRESOLVED (genuine conflicting type at the same/adjacent tile position):\n"
        "- Ref#/Tag: <value> | Candidates: <type A> or <type B> | Tile: <tile>\n"
    )

    # --- Non-piping-diagram document passes (data sheets, isometrics/GA, matrices) ---
    TABLE_FIELD_PROMPT = (
        "You are reading a scanned engineering document that is primarily tabular/field-based "
        "(e.g. a data sheet, equipment sizing report, proposal, or spec sheet), not a piping "
        "diagram.\n\n"
        "Extract every labelled field and every table exactly as written.\n"
        "RULES:\n"
        "- For label:value fields (e.g. 'Design Pressure: 150 psig'), extract as Label = Value\n"
        "- For tables, reproduce each row with its column headers\n"
        "- Preserve units exactly as written\n"
        "- Do not calculate, convert, or infer values not explicitly present\n"
        "- If a field or cell is blank/illegible, write N/A\n\n"
        "OUTPUT FORMAT:\n"
        "FIELDS:\n"
        "- <Label>: <Value>\n\n"
        "TABLES:\n"
        "<table name/title if any>\n"
        "| <header1> | <header2> | ... |\n"
        "| <row values> |\n\n"
        "(repeat for each table found)\n"
    )

    DIMENSION_CALLOUT_PROMPT = (
        "You are reading a scanned isometric or general arrangement (GA) drawing.\n\n"
        "Extract every dimension, weld number, coordinate, elevation, and callout visible.\n"
        "RULES:\n"
        "- Dimensions appear as numbers with arrows or extension lines -- record value and units\n"
        "- Weld numbers are small circled or boxed numbers along a pipe run -- list each with its "
        "approximate location description (e.g. 'near flange at north end')\n"
        "- Elevations/coordinates are typically labelled EL. or N/E/W/S with a numeric value\n"
        "- Note pipe spec labels and line numbers if present\n"
        "- Do not invent values not explicitly shown\n\n"
        "OUTPUT FORMAT:\n"
        "DIMENSIONS:\n- <value> <units>: <what it measures>\n\n"
        "WELD NUMBERS:\n- <number>: <location description>\n\n"
        "ELEVATIONS / COORDINATES:\n- <label>: <value>\n\n"
        "OTHER CALLOUTS:\n- <text as written>\n"
    )

    MATRIX_PROMPT = (
        "You are reading a scanned cause-and-effect matrix or HAZOP worksheet.\n\n"
        "RULES:\n"
        "- Extract row labels (causes/deviations) and column labels (effects/actions) exactly as "
        "written\n"
        "- For each marked intersection (X, dot, or shaded cell), record which row/column it links\n"
        "- Preserve any severity/likelihood ratings shown\n"
        "- Do not infer a mark where none is visibly present\n\n"
        "OUTPUT FORMAT:\n"
        "ROWS:\n- <row label>\n\n"
        "COLUMNS:\n- <column label>\n\n"
        "MARKED INTERSECTIONS:\n- <row label> -> <column label>\n"
    )

    # Maps a classifier/AI doc_type guess to which pass category should run.
    # Anything not listed (including "Unknown") defaults to the full diagram pipeline,
    # matching the existing coordinator philosophy: unclassified process content should
    # still get the thorough pass rather than a shallow one.
    _ISO_GA_TYPES = {"Isometric", "GA Drawing"}
    _TABULAR_TYPES = {"Data Sheet", "Equipment Sizing", "Proposal", "Report", "Heat & Energy Balance"}
    _MATRIX_TYPES = {"Cause & Effect", "Hazop"}

    def extract_text_sample(self, path: Path, max_chars: int = 2000) -> str:
        ext = path.suffix.lower()
        try:
            if ext == ".pdf":
                return self._read_pdf(path, max_chars)
            elif ext in (".docx", ".doc"):
                return self._read_docx(path, max_chars)
            elif ext in (".xlsx", ".xls"):
                return self._read_xlsx(path, max_chars)
            elif ext in (".pptx", ".ppt"):
                return self._read_pptx(path, max_chars)
            elif ext == ".txt":
                return path.read_text(errors="ignore")[:max_chars]
            elif ext in (".png", ".jpg", ".jpeg"):
                return ""
        except Exception:
            return ""
        return ""

    def extract_vision_description(self, path: Path, api_key: str = "", doc_type_hint: str = "") -> str:
        ext = path.suffix.lower()
        try:
            stat = path.stat()
            cache_key = (str(path.resolve()), stat.st_mtime_ns, stat.st_size, (doc_type_hint or "").strip())
            cached = self._vision_cache.get(cache_key)
            if cached is not None:
                return cached
        except Exception:
            cache_key = None

        try:
            if ext in (".png", ".jpg", ".jpeg"):
                image_b64 = base64.b64encode(path.read_bytes()).decode("utf-8")
                mime = "image/png" if ext == ".png" else "image/jpeg"
                pil_image = self._b64_to_pil(image_b64)
            elif ext == ".pdf":
                pil_image = self._pdf_to_pil(path)
                image_b64 = None
                mime = "image/png"
            else:
                return ""
        except Exception as e:
            return f"Image load failed: {e}"

        result = self._multi_pass_analysis(pil_image, mime, api_key, doc_type_hint=doc_type_hint)
        if cache_key is not None and result and not result.startswith("Image load failed"):
            self._vision_cache[cache_key] = result
        return result

    def _categorize_doc_type(self, doc_type_hint: str) -> str:
        """Which pass set applies. Unrecognised/Unknown types default to the full
        diagram pipeline rather than a shallow one -- an unclassified drawing is more
        likely under-classified than genuinely non-diagrammatic."""
        dt = (doc_type_hint or "").strip()
        if dt in self._ISO_GA_TYPES:
            return "iso_ga"
        if dt in self._MATRIX_TYPES:
            return "matrix"
        if dt in self._TABULAR_TYPES:
            return "tabular"
        return "diagram"

    def _doc_type_from_context(self, context_text: str) -> str:
        """Recover the model-read document type from the context pass when the
        filename/rule classifier could not identify a scanned drawing beforehand."""
        m = re.search(r"(?im)^\s*(?:\d+[.)]\s*)?(?:[-*]\s*)?\**DOCUMENT TYPE\**\s*:\s*([^\n]+)", context_text or "")
        if not m:
            return ""
        value = m.group(1).strip().lower()
        aliases = {
            "p&id": "P&ID",
            "pid": "P&ID",
            "pfd": "PFD",
            "process flow diagram": "PFD",
            "system diagram": "System Diagram",
            "data sheet": "Data Sheet",
            "datasheet": "Data Sheet",
            "isometric": "Isometric",
            "ga": "GA Drawing",
            "ga drawing": "GA Drawing",
            "general arrangement": "GA Drawing",
            "cause & effect": "Cause & Effect",
            "cause and effect": "Cause & Effect",
            "hazop": "Hazop",
        }
        for key, canonical in aliases.items():
            if key in value:
                return canonical
        return ""

    def _adaptive_grid(self, width: int, height: int, target_edge: int | None = None) -> tuple[int, int]:
        """Compute the first-pass tile grid. The target balances crop coverage against
        small-symbol readability; hard cases are re-cropped at higher magnification by the
        targeted verification passes rather than forcing every region through a finer grid.
        """
        target_edge = target_edge or self.BASE_TILE_TARGET
        cols = max(1, -(-width // target_edge))    # ceil division
        rows = max(1, -(-height // target_edge))
        return cols, rows

    def _needs_zoom_verification(self, tile_result: str) -> bool:
        if not tile_result:
            return False
        lower = tile_result.lower()
        return any(t in lower for t in ("unknown", "uncertain", "candidates:"))

    def _is_dense_tile(self, tile_result: str, threshold: int = 4) -> bool:
        """Dense tiles (many items packed into one crop) are exactly where LLM vision
        undercounts and conflates nearby items -- e.g. missing a fitting when three others
        sit close together, or blending two labels. This doesn't require the model to admit
        uncertainty; it's just a property of how much the tile is reporting."""
        if not tile_result:
            return False
        item_lines = [l for l in tile_result.splitlines() if l.strip().startswith("-")]
        return len(item_lines) >= threshold

    def _instrument_needs_zoom(self, tile_result: str) -> bool:
        """Only pay for a second instrument look when the first pass admits uncertainty.
        This keeps easy tiles single-pass while preserving the high-magnification escape
        hatch for the tiny/blurred tags that actually need it.
        """
        if not tile_result:
            return False
        lower = tile_result.lower()
        return "unreadable" in lower or "????" in tile_result or "?" in tile_result

    def _reconcile_pipe_specs(self, raw_specs: list[str], reference_spec: str, api_key: str) -> str:
        if not raw_specs:
            return "CONFIRMED SPECS:\n(none found)"
        key = api_key or os.environ.get("OPENAI_API_KEY", "")
        if not key:
            seen = []
            for s in raw_specs:
                if s not in seen:
                    seen.append(s)
            return "CONFIRMED SPECS:\n" + "\n".join(f"- {s}" for s in seen)
        from openai import OpenAI
        client = OpenAI(api_key=key)
        try:
            response = client.chat.completions.create(
                model="gpt-4o",
                max_tokens=800,
                messages=[{
                    "role": "user",
                    "content": self.PIPE_SPEC_RECONCILE_PROMPT.format(
                        reference_spec=reference_spec or "none found",
                        raw_specs="\n".join(raw_specs),
                    ),
                }],
            )
            self._track_usage("pipe_spec_reconcile", response.usage)
            return response.choices[0].message.content or ""
        except Exception as e:
            return "CONFIRMED SPECS:\n" + "\n".join(f"- {s}" for s in raw_specs) + f"\n\n[Reconciliation failed: {e}]"

    def _dedupe_reconciled_valves(self, text: str) -> str:
        """Deterministic safety net run AFTER LLM reconciliation (or in place of it when no
        API key is available). Groups entries by ref#/tag, then clusters each group by tile
        adjacency: entries whose tiles are the same or adjacent get merged (a real duplicate
        seen twice via overlapping crops); entries whose tiles are far apart stay as separate
        output entries (a reference number legitimately reused across multiple distinct
        physical items, e.g. a 'TYP.' callout). This is what prevents the dedup net itself
        from re-introducing the exact bug it exists to fix -- collapsing distinct real items
        into one false 'conflicting reads' entry just because they share a label. Lines that
        don't match the expected format are passed through unchanged.
        """
        if not text or not text.strip():
            return text

        header_pattern = re.compile(r"^\[Tile\s+(R\d+C\d+)")
        entry_pattern = re.compile(
            r"^-\s*(?P<type>[^|]+?)\s*\|\s*Line:\s*(?P<line>[^|]*?)\s*\|\s*Tag:\s*(?P<tag>[^|]*?)\s*\|\s*"
            r"Ref#:\s*(?P<ref>[^|]*?)(?:\s*\|\s*Tile:\s*(?P<tile>R\d+C\d+))?\s*$"
        )
        unresolved_pattern = re.compile(
            r"^-\s*Ref#/Tag:\s*(?P<key>[^|]+?)\s*\|\s*Candidates:\s*(?P<candidates>[^|]*?)"
            r"(?:\s*\|\s*Tile:\s*(?P<tile>R\d+C\d+))?\s*$"
        )

        raw_entries: list[dict] = []
        passthrough: list[str] = []
        current_tile = None

        for raw_line in text.splitlines():
            line = raw_line.strip()
            hm = header_pattern.match(line)
            if hm:
                current_tile = hm.group(1)
                continue

            if not line.startswith("-") or not line.strip("- "):
                if line:
                    passthrough.append(raw_line)
                continue

            m = entry_pattern.match(line)
            if m:
                tag = (m.group("tag") or "").strip().rstrip(".")
                ref = (m.group("ref") or "").strip().rstrip(".")
                if tag and tag.lower() not in ("none", "n/a", "-", ""):
                    dedupe_key = f"tag:{tag.lower()}"
                elif ref and ref.lower() not in ("none", "n/a", "-", ""):
                    dedupe_key = f"ref:{ref.lower()}"
                else:
                    passthrough.append(raw_line)
                    continue
                raw_entries.append({
                    "key": dedupe_key,
                    "types": [t.strip() for t in re.split(r"\bor\b", m.group("type").strip()) if t.strip()],
                    "line": m.group("line").strip(),
                    "tag": tag, "ref": ref,
                    "tile": m.group("tile") or current_tile or "",
                })
                continue

            um = unresolved_pattern.match(line)
            if um:
                key_raw = um.group("key").strip()
                dedupe_key = f"ref:{key_raw.lower()}" if key_raw.isdigit() else f"tag:{key_raw.lower()}"
                raw_entries.append({
                    "key": dedupe_key,
                    "types": [c.strip() for c in re.split(r"\bor\b", um.group("candidates")) if c.strip()],
                    "line": "",
                    "tag": key_raw if not key_raw.isdigit() else "",
                    "ref": key_raw if key_raw.isdigit() else "",
                    "tile": um.group("tile") or current_tile or "",
                })
                continue

            if line:
                passthrough.append(raw_line)

        if not raw_entries:
            return text

        by_key: dict[str, list[dict]] = {}
        key_order: list[str] = []
        for e in raw_entries:
            if e["key"] not in by_key:
                by_key[e["key"]] = []
                key_order.append(e["key"])
            by_key[e["key"]].append(e)

        rebuilt = []
        for key in key_order:
            clusters: list[list[dict]] = []
            for e in by_key[key]:
                placed = False
                for cluster in clusters:
                    if any(self._tiles_adjacent(e["tile"], m["tile"]) for m in cluster):
                        cluster.append(e)
                        placed = True
                        break
                if not placed:
                    clusters.append([e])

            for cluster in clusters:
                types: list[str] = []
                for e in cluster:
                    for t in e["types"]:
                        if t not in types:
                            types.append(t)
                line_spec = next(
                    (e["line"] for e in cluster if e["line"] and e["line"].lower() not in ("none", "n/a", "-")),
                    "not stated",
                )
                tag = next((e["tag"] for e in cluster if e["tag"]), "none")
                ref = next((e["ref"] for e in cluster if e["ref"]), "none")
                tiles_seen = sorted(set(e["tile"] for e in cluster if e["tile"]))
                tile_str = "/".join(tiles_seen) if tiles_seen else "unknown"
                type_str = " or ".join(types) if types else "UNKNOWN"
                suffix = "  (conflicting reads at this location -- verify against drawing)" if len(types) > 1 else ""
                rebuilt.append(
                    f"- {type_str} | Line: {line_spec} | Tag: {tag} | Ref#: {ref} | Tile: {tile_str}{suffix}"
                )

        result_lines = rebuilt + [p for p in passthrough if p.strip()]
        return "\n".join(result_lines)

    def _reconcile_valve_survey(self, raw_blocks: list[str], instrument_tags: list[str], api_key: str) -> str:
        if not raw_blocks:
            return "(none found)"
        key = api_key or os.environ.get("OPENAI_API_KEY", "")
        if not key:
            return self._dedupe_reconciled_valves("\n".join(raw_blocks))
        from openai import OpenAI
        client = OpenAI(api_key=key)
        try:
            response = client.chat.completions.create(
                model="gpt-4o",
                max_tokens=1500,
                messages=[{
                    "role": "user",
                    "content": self.VALVE_RECONCILE_PROMPT.format(
                        instrument_tags="\n".join(instrument_tags) if instrument_tags else "none",
                        raw_survey="\n---\n".join(raw_blocks),
                    ),
                }],
            )
            result = response.choices[0].message.content or ""
            self._track_usage("valve_reconcile", response.usage)
            return self._dedupe_reconciled_valves(result)
        except Exception as e:
            return self._dedupe_reconciled_valves("\n".join(raw_blocks)) + f"\n\n[Reconciliation failed: {e}]"

    def _multi_pass_analysis(self, pil_image, mime: str, api_key: str, doc_type_hint: str = "") -> str:
        width, height = pil_image.size
        category = self._categorize_doc_type(doc_type_hint)

        # --- Pass 1: Context (full image, layout / title block / equipment) -- always runs ---
        context_b64 = self._pil_to_b64(pil_image)
        context_text = self._call_vision(context_b64, mime, self.CONTEXT_PROMPT, api_key, max_tokens=1200, label="context")
        description = "=== CONTEXT ANALYSIS ===\n" + context_text

        # If the rule-based classifier had no usable hint, the already-paid context pass
        # can route the remaining work. This avoids running a full P&ID scan on a scanned
        # data sheet/GA/matrix simply because its filename was uninformative.
        if not doc_type_hint or doc_type_hint.strip() == "Unknown":
            context_doc_type = self._doc_type_from_context(context_text)
            if context_doc_type:
                category = self._categorize_doc_type(context_doc_type)

        if category == "tabular":
            table_b64 = self._pil_to_b64(pil_image)
            table_text = self._call_vision(table_b64, mime, self.TABLE_FIELD_PROMPT, api_key, max_tokens=1500, label="table_field")
            description += "\n\n=== FIELDS AND TABLES ===\n" + table_text
            return description

        if category == "matrix":
            matrix_b64 = self._pil_to_b64(pil_image)
            matrix_text = self._call_vision(matrix_b64, mime, self.MATRIX_PROMPT, api_key, max_tokens=1500, label="matrix")
            description += "\n\n=== MATRIX CONTENTS ===\n" + matrix_text
            return description

        if category == "iso_ga":
            diagram_body = pil_image.crop((0, int(height * 0.03), width, int(height * 0.9)))
            cols, rows = self._adaptive_grid(diagram_body.size[0], diagram_body.size[1])
            callout_results = []
            for tile in self._make_tiles(diagram_body, cols=cols, rows=rows, overlap_frac=0.12):
                tile_b64 = self._pil_to_b64(tile)
                result = self._call_vision(tile_b64, mime, self.DIMENSION_CALLOUT_PROMPT, api_key, max_tokens=500, label="dimension_callout")
                if result.strip():
                    callout_results.append(result.strip())
            description += "\n\n=== DIMENSIONS AND CALLOUTS ===\n" + "\n---\n".join(callout_results)

            spec_results = []
            spec_cols = max(4, -(-width // self.BASE_TILE_TARGET))  # single-row strips: width-only
            for tile in self._make_tiles(diagram_body, cols=spec_cols, rows=1, overlap_frac=0.10):
                tile_b64 = self._pil_to_b64(tile)
                result = self._call_vision(tile_b64, mime, self.PIPE_SPEC_PROMPT, api_key, max_tokens=300, label="iso_pipe_spec_tile")
                if result.strip():
                    spec_results.append(result.strip())
            description += "\n\n=== PIPE SPECIFICATIONS ===\n" + "\n".join(spec_results)
            return description

        # --- category == "diagram": P&ID / System Diagram / PFD / Unknown -- full pipeline ---

        # --- Pass 2: Instrument tags -- single adaptive-resolution pass ---
        all_raw_tags = []
        i_cols, i_rows = self._adaptive_grid(width, height)
        for tile in self._make_tiles(pil_image, cols=i_cols, rows=i_rows, overlap_frac=0.12):
            tile_b64 = self._pil_to_b64(tile)
            result = self._call_vision(tile_b64, mime, self.TILE_INSTRUMENT_PROMPT, api_key, max_tokens=450, label="instrument_tile")
            all_raw_tags.extend(self._parse_tag_list(result))

            # If the model explicitly says a tag is unreadable, split only that tile and
            # inspect again at higher magnification.
            if self._instrument_needs_zoom(result):
                for sub in self._make_tiles(tile, cols=2, rows=2, overlap_frac=0.15):
                    sub_b64 = self._pil_to_b64(sub)
                    zoom_out = self._call_vision(
                        sub_b64, mime, self.TILE_INSTRUMENT_PROMPT, api_key, max_tokens=300, label="instrument_zoom"
                    )
                    all_raw_tags.extend(self._parse_tag_list(zoom_out))

        reconciled_tags = self._validate_tags_local(all_raw_tags)

        # --- Pass 3: Pipe spec reading -- horizontal strips across diagram body ---
        pipe_spec_results = []
        pipe_spec_zoom_results = []
        diagram_body = pil_image.crop((0, int(height * 0.05), width, int(height * 0.82)))
        spec_cols = max(4, -(-diagram_body.size[0] // self.BASE_TILE_TARGET))  # single-row strips: width-only
        for tile in self._make_tiles(diagram_body, cols=spec_cols, rows=1, overlap_frac=0.10):
            tile_b64 = self._pil_to_b64(tile)
            result = self._call_vision(tile_b64, mime, self.PIPE_SPEC_PROMPT, api_key, max_tokens=300, label="pipe_spec_tile")
            if result.strip():
                pipe_spec_results.append(result.strip())

                # Targeted zoom: fires on an explicit '?' (model flagged a char as unclear)
                # or a dense strip (many labels in one crop -- exactly where digits from
                # adjacent labels bleed into each other).
                if "?" in result or self._is_dense_tile(result, threshold=3):
                    for sub in self._make_tiles(tile, cols=2, rows=1, overlap_frac=0.2):
                        sub_b64 = self._pil_to_b64(sub)
                        zoom_out = self._call_vision(
                            sub_b64, mime, self.PIPE_SPEC_DETAIL_ZOOM_PROMPT, api_key, max_tokens=250, label="pipe_spec_zoom"
                        )
                        if zoom_out.strip():
                            pipe_spec_zoom_results.append(zoom_out.strip())

        raw_specs = []
        reference_spec = ""
        for block in pipe_spec_results + pipe_spec_zoom_results:
            for line in block.splitlines():
                spec = line.strip().strip("-").strip()
                if not spec:
                    continue
                if spec.upper().startswith("TYP:"):
                    reference_spec = spec.split(":", 1)[1].strip()
                else:
                    raw_specs.append(spec)

        reconciled_specs = self._reconcile_pipe_specs(raw_specs, reference_spec, api_key)
        pipe_spec_section = "\n=== PIPE SPECIFICATIONS (reconciled across tiles) ===\n" + reconciled_specs

        # --- Pass 4: Valve survey -- adaptive grid over the diagram body, every tile ---
        valve_results = []
        zoom_results = []
        v_cols, v_rows = self._adaptive_grid(diagram_body.size[0], diagram_body.size[1])
        for row, col, tile in self._make_tiles_indexed(diagram_body, cols=v_cols, rows=v_rows, overlap_frac=0.15):
            tile_label = f"R{row}C{col}"
            tile_b64 = self._pil_to_b64(tile)
            result = self._call_vision(tile_b64, mime, self.VALVE_SURVEY_PROMPT, api_key, max_tokens=550, label="valve_tile")
            if result.strip() and len(result.strip()) > 20:
                valve_results.append(f"[Tile {tile_label}]\n{result.strip()}")

                # Re-examine THIS tile at higher magnification if the first pass flagged
                # ambiguity OR the tile is dense (a silent miss never self-reports as
                # "unknown" -- density is the only signal that catches it).
                if self._needs_zoom_verification(result) or self._is_dense_tile(result, threshold=6):
                    for sub in self._make_tiles(tile, cols=2, rows=2, overlap_frac=0.2):
                        sub_b64 = self._pil_to_b64(sub)
                        zoom_out = self._call_vision(
                            sub_b64, mime, self.VALVE_DETAIL_ZOOM_PROMPT, api_key, max_tokens=500, label="valve_zoom"
                        )
                        if zoom_out.strip() and len(zoom_out.strip()) > 15:
                            # Zoom subdivides within the SAME parent tile, so it shares that
                            # tile's position label -- it's a closer look at the same region,
                            # not a different location.
                            zoom_results.append(f"[Tile {tile_label}, zoom verification]\n{zoom_out.strip()}")

        reconciled_valves = self._reconcile_valve_survey(
            valve_results + zoom_results,
            reconciled_tags,
            api_key,
        )
        valve_section = "\n\n=== VALVE AND FITTING SURVEY (reconciled across tiles) ===\n" + reconciled_valves

        description += (
            "\n\n=== INSTRUMENT TAGS (multi-tile extraction) ===\n"
            + "\n".join(reconciled_tags)
            + pipe_spec_section
            + valve_section
        )

        return description

    def _make_tiles(self, image, cols: int, rows: int, overlap_frac: float = 0.1):
        from PIL import Image
        w, h = image.size
        tile_w = int(w / cols)
        tile_h = int(h / rows)
        overlap_x = int(tile_w * overlap_frac)
        overlap_y = int(tile_h * overlap_frac)
        tiles = []
        for row in range(rows):
            for col in range(cols):
                x0 = max(0, col * tile_w - overlap_x)
                y0 = max(0, row * tile_h - overlap_y)
                x1 = min(w, (col + 1) * tile_w + overlap_x)
                y1 = min(h, (row + 1) * tile_h + overlap_y)
                tiles.append(image.crop((x0, y0, x1, y1)))
        return tiles

    def _make_tiles_indexed(self, image, cols: int, rows: int, overlap_frac: float = 0.1):
        """Same crop geometry as _make_tiles, but yields (row, col, tile) so callers can
        stamp each reading with WHERE it came from. This costs nothing extra to compute --
        the grid position is already known in code, no need to ask the model for it."""
        w, h = image.size
        tile_w = int(w / cols)
        tile_h = int(h / rows)
        overlap_x = int(tile_w * overlap_frac)
        overlap_y = int(tile_h * overlap_frac)
        out = []
        for row in range(rows):
            for col in range(cols):
                x0 = max(0, col * tile_w - overlap_x)
                y0 = max(0, row * tile_h - overlap_y)
                x1 = min(w, (col + 1) * tile_w + overlap_x)
                y1 = min(h, (row + 1) * tile_h + overlap_y)
                out.append((row, col, image.crop((x0, y0, x1, y1))))
        return out

    @staticmethod
    def _tiles_adjacent(tile_a: str, tile_b: str) -> bool:
        """Two tile labels like 'R2C3' -- adjacent (including same) means overlapping crops
        could plausibly have captured the same physical symbol twice."""
        import re as _re
        ma = _re.match(r"R(\d+)C(\d+)", tile_a or "")
        mb = _re.match(r"R(\d+)C(\d+)", tile_b or "")
        if not ma or not mb:
            return False
        r1, c1 = int(ma.group(1)), int(ma.group(2))
        r2, c2 = int(mb.group(1)), int(mb.group(2))
        return abs(r1 - r2) <= 1 and abs(c1 - c2) <= 1

    def _pil_to_b64(self, image) -> str:
        buf = io.BytesIO()
        image.save(buf, format="PNG")
        return base64.b64encode(buf.getvalue()).decode("utf-8")

    def _b64_to_pil(self, b64_str: str):
        from PIL import Image
        data = base64.b64decode(b64_str)
        return Image.open(io.BytesIO(data))

    def _pdf_to_pil(self, path: Path):
        import fitz
        from PIL import Image
        doc = fitz.open(str(path))
        page = doc[0]
        mat = fitz.Matrix(4.0, 4.0)
        pix = page.get_pixmap(matrix=mat)
        doc.close()
        img_data = pix.tobytes("png")
        return Image.open(io.BytesIO(img_data))

    def _call_vision(self, image_b64: str, mime: str, prompt: str, api_key: str, max_tokens: int = 800, label: str = "vision") -> str:
        from openai import OpenAI
        key = api_key or os.environ.get("OPENAI_API_KEY", "")
        if not key:
            return ""
        client = OpenAI(api_key=key)
        try:
            response = client.chat.completions.create(
                model="gpt-4o",
                max_tokens=max_tokens,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime};base64,{image_b64}",
                                    "detail": "high",
                                },
                            },
                        ],
                    }
                ],
            )
            self._track_usage(label, response.usage)
            return response.choices[0].message.content or ""
        except Exception as e:
            return f"Vision call failed: {e}"

    def _parse_tag_list(self, text: str) -> list[str]:
        tags = []
        for line in text.splitlines():
            line = line.strip().strip("-").strip("*").strip()
            if not line or line.startswith("[") or line.lower().startswith("unreadable"):
                continue
            match = re.match(r"([A-Z]{2,5}-\d{3,6}[A-Z]?)(\s.*)?$", line)
            if match:
                tags.append(match.group(1))
        return tags

    def _validate_tags_local(self, tags: list[str]) -> list[str]:
        validated = []
        seen = set()
        for tag in tags:
            m = re.match(r"^([A-Z]+)-(\d+)([A-Z]?)$", tag)
            if not m:
                continue
            prefix = m.group(1)
            if prefix not in self.VALID_ISA_PREFIXES:
                continue
            if tag not in seen:
                seen.add(tag)
                validated.append(tag)
        return validated

    def compact_vision_description(self, vision_text: str, section_names: tuple[str, ...] | None = None) -> str:
        """Return only reconciled/structured vision sections for downstream LLM calls.

        Older runs may still contain RAW TILE OUTPUTS; those are deliberately excluded
        because they duplicate the reconciled instrument list and amplify downstream tokens.
        """
        if not vision_text:
            return ""
        sections: dict[str, str] = {}
        current = None
        buf: list[str] = []
        for line in vision_text.splitlines():
            m = re.match(r"^===\s*(.+?)\s*===$", line.strip())
            if m:
                if current is not None:
                    sections[current] = "\n".join(buf).strip()
                current = m.group(1).strip()
                buf = []
            elif current is not None:
                buf.append(line)
        if current is not None:
            sections[current] = "\n".join(buf).strip()

        wanted = section_names or tuple(k for k in sections if not k.startswith("RAW TILE OUTPUTS"))
        chunks = []
        for name in wanted:
            body = sections.get(name, "").strip()
            if body:
                chunks.append(f"=== {name} ===\n{body}")
        return "\n\n".join(chunks) if chunks else vision_text

    def extract_metadata(self, path: Path, text_sample: str, doc_type: str) -> dict:
        meta = {
            "doc_type": doc_type,
            "project_number": self._find_project_number(path.name, text_sample),
            "revision": self._find_revision(path.name, text_sample),
            "description": self._find_description(path.name),
            "client": self._find_client(text_sample),
            "date": self._find_date(text_sample),
            "unit_operations": self._find_unit_operations(text_sample),
            "instrumentation": self._find_instrumentation(text_sample),
        }
        return {k: v for k, v in meta.items() if v}

    def extract_metadata_from_vision(self, vision_text: str, path: Path) -> dict:
        instrument_section = ""
        if "=== INSTRUMENT TAGS" in vision_text:
            parts = vision_text.split("=== INSTRUMENT TAGS (multi-tile extraction) ===")
            if len(parts) > 1:
                instrument_section = parts[1].split("===")[0].strip()

        context_section = ""
        if "=== CONTEXT ANALYSIS ===" in vision_text:
            parts = vision_text.split("=== CONTEXT ANALYSIS ===")
            if len(parts) > 1:
                context_section = parts[1].split("===")[0].strip()

        combined = context_section + " " + instrument_section

        tags = []
        for line in instrument_section.splitlines():
            line = line.strip()
            m = re.match(r"([A-Z]{2,5}-\d{3,6}[A-Z]?\??)", line)
            if m:
                tags.append(m.group(1))

        meta = {
            "project_number": self._find_project_number(path.name, context_section),
            "revision": self._find_revision(path.name, context_section),
            "description": self._find_description(path.name),
            "client": self._find_client(context_section),
            "date": self._find_date(context_section),
            "unit_operations": self._find_unit_operations(context_section),
            "instrumentation": tags if tags else self._find_instrumentation(combined),
            "vision_description": vision_text,
        }
        return {k: v for k, v in meta.items() if v}

    def _find_project_number(self, filename: str, text: str) -> str:
        combined = filename + " " + text

        explicit_patterns = [
            (r"Project\s*(?:No|Number|#|Num)[.:\s]+([A-Z]{2,4}[-_]?\d{2,6})", 1),
            (r"\b(PRJ[-_]?\d{2,6})\b", 1),
            (r"\b(JOB[-_]?\d{2,6})\b", 1),
        ]
        for pat, group in explicit_patterns:
            m = re.search(pat, combined, re.IGNORECASE)
            if m:
                candidate = m.group(group).strip()
                if re.match(r"^\d{7,}$", candidate):
                    continue
                return candidate

        return ""

    def _find_revision(self, filename: str, text: str) -> str:
        patterns = [
            r"\bRev\.?\s*[A-Z0-9]{1,3}\b",
            r"\bR[0-9]{1,2}\b",
            r"\b[Rr]evision\s+[A-Z0-9]{1,3}\b",
        ]
        combined = filename + " " + text
        for pat in patterns:
            m = re.search(pat, combined, re.IGNORECASE)
            if m:
                return m.group(0).strip()
        return ""

    def _find_description(self, filename: str) -> str:
        stem = Path(filename).stem
        stem = re.sub(r"(PRJ[-_]?\d+|Rev\.?\s*\w+|R\d+|\d{4,})", "", stem, flags=re.IGNORECASE)
        stem = re.sub(r"[-_]+", " ", stem).strip()
        words = [w for w in stem.split() if len(w) > 1]
        return " ".join(words[:6])

    def _find_client(self, text: str) -> str:
        patterns = [
            r"Client[:\s]+([A-Za-z0-9 &,.\-]{3,40})",
            r"Prepared\s+for[:\s]+([A-Za-z0-9 &,.\-]{3,40})",
        ]
        for pat in patterns:
            m = re.search(pat, text, re.IGNORECASE)
            if m:
                return m.group(1).strip()
        return ""

    def _find_date(self, text: str) -> str:
        patterns = [
            r"\b(\d{1,2}[-/]\d{1,2}[-/]\d{2,4})\b",
            r"\b(\d{1,2}\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4})\b",
            r"\b((?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4})\b",
        ]
        for pat in patterns:
            m = re.search(pat, text, re.IGNORECASE)
            if m:
                return m.group(1).strip()
        return ""

    def _find_unit_operations(self, text: str) -> list[str]:
        keywords = [
            "separator", "heat exchanger", "compressor", "pump", "vessel",
            "column", "distillation", "absorber", "stripper", "reactor",
            "filter", "scrubber", "cooler", "heater", "reboiler", "condenser",
            "flash drum", "knock-out drum", "slug catcher", "electric motor",
            "driver motor", "aftercooler", "suction drum", "discharge drum",
            "control valve", "relief valve", "check valve", "blowdown",
            "gas seal", "seal panel", "skid",
        ]
        found = []
        lower = text.lower()
        for kw in keywords:
            if kw in lower and kw not in found:
                found.append(kw)
        return found

    def _find_instrumentation(self, text: str) -> list[str]:
        candidates = re.findall(r"\b([A-Z]{2,5}-\d{3,5}[A-Z]?)\b", text)
        seen = []
        for tag in candidates:
            prefix = re.match(r"^([A-Z]+)", tag)
            if not prefix:
                continue
            if prefix.group(1) not in self.VALID_ISA_PREFIXES:
                continue
            if tag not in seen:
                seen.append(tag)
        return seen[:60]

    def _read_pdf(self, path: Path, max_chars: int) -> str:
        import fitz
        doc = fitz.open(str(path))
        text = ""
        for page in doc:
            text += page.get_text()
            if len(text) >= max_chars:
                break
        doc.close()
        return text[:max_chars]

    def _read_docx(self, path: Path, max_chars: int) -> str:
        from docx import Document
        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs)
        return text[:max_chars]

    def _read_xlsx(self, path: Path, max_chars: int) -> str:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        text = ""
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                text += " ".join(str(c) for c in row if c is not None) + "\n"
                if len(text) >= max_chars:
                    break
        return text[:max_chars]

    def _read_pptx(self, path: Path, max_chars: int) -> str:
        from pptx import Presentation
        prs = Presentation(str(path))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                if len(text) >= max_chars:
                    break
        return text[:max_chars]
